from pptx import Presentation
from wand import assertions
from wand.image import Image
from wand.display import display
import os
from pptx.util import Inches

prs = Presentation()
test_image=[
  "/home/manish/Downloads/cc1.jpg",
  "/home/manish/Downloads/cc3.jpg",
  "/home/manish/Downloads/cc4.jpg",
  "/home/manish/Downloads/cc5.jpg",
  "/home/manish/Downloads/cc6.jpg"]


output_img=[
  "img1.jpeg", 
  "img2.jpeg",
  "img3.jpeg", 
  "img4.jpeg",
  "img5.jpeg" ]

b=["jpeg", "jpg"]
for i in range(len(test_image)):
  assert test_image[i][-3:] in b, "oops image not in supported format"
  with Image(filename =test_image[i]) as image:
      # Import the watermark image
      with Image(filename ='/home/manish/Downloads/cc2.png' ) as water_mark:
          # Clone the image in order to process
          water_mark.resize(1300,500)
          with image.clone() as watermark:
              # Invoke watermark function with [watermark image, transparency as 0
              # left as 10 and top as 180]
              watermark.watermark(water_mark, 0, 10, 180)
                  # Save the image
              display(watermark)
              watermark.save(filename =output_img[i])
def fun1(data):
  layout = prs.slide_layouts[data[3]]
  slide=prs.slides.add_slide(layout)
  title=slide.shapes.title
  title.text=data[0]
  subtitle=slide.placeholders[1]
  subtitle.text=data[1]
  top = Inches(2.3) 
  left = Inches(0.1) 
  height = Inches(5) 
  slide.shapes.add_picture(data[2],left,top, height=height) 
     

slides = [
    ["It's me Title1",       #data[0]
     "It's me sub Title1)",
     output_img[0],
     3],
    ["It's me Title2",       #data[1]
     "It's me sub Title2)",
     output_img[1],
     3],
    ["It's me Title3",       #data[2]
     "It's me sub Title3",
     output_img[2],
     3],
    ["It's me Title4",       #data[3]
     "It's me sub Title4",
     output_img[3],
     3],
    ["It's me Title5",       #data[4]
     "It's me sub Title5",
     output_img[4],
     3]
    ]
 
for each_slide in slides:
    fun1(each_slide)
 
prs.save("output.pptx")
#os.startfile("stack.pptx")  "this command run only windows os"
